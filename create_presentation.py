import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Palette (Dark Mode Glassmorphism Theme)
    BG_COLOR = RGBColor(15, 23, 42)        # Deep Slate / Navy
    CARD_BG = RGBColor(30, 41, 59)         # Dark Card Slate
    ACCENT_CYAN = RGBColor(6, 182, 212)    # Cyan Accent
    ACCENT_PURPLE = RGBColor(168, 85, 247) # Purple Accent
    TEXT_WHITE = RGBColor(248, 250, 252)   # Bright White
    TEXT_MUTED = RGBColor(148, 163, 184)  # Muted Gray

    def add_background(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="COURSE CRAFT AI • SYIT PROJECT PRESENTATION"):
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.06)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_CYAN
        accent_line.line.fill.background()

        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.733), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1)

    card = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ACCENT_CYAN
    card.line.width = Pt(1.5)

    tb = slide1.shapes.add_textbox(Inches(2.0), Inches(1.8), Inches(9.333), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🎓 CourseCraft AI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Intelligent Python Course Generator Platform"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\nSYIT College Final Project Presentation"
    p3.font.size = Pt(16)
    p3.font.color.rgb = ACCENT_PURPLE
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "\nPresented by: Munim Gupta  |  Technology: Python, Flask, SQLite/PostgreSQL, Web Engine"
    p4.font.size = Pt(14)
    p4.font.color.rgb = TEXT_MUTED
    p4.alignment = PP_ALIGN.CENTER

    p5 = tf.add_paragraph()
    p5.text = "Live App: course-craft-ai-xlxe.onrender.com"
    p5.font.size = Pt(13)
    p5.font.color.rgb = ACCENT_CYAN
    p5.alignment = PP_ALIGN.CENTER

    # SLIDE 2: Problem Statement & Solution
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2)
    add_header(slide2, "Problem Statement & Proposed Solution")

    prob_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8))
    prob_box.fill.solid()
    prob_box.fill.fore_color.rgb = CARD_BG
    prob_box.line.color.rgb = RGBColor(239, 68, 68)
    prob_box.line.width = Pt(1.5)

    tf_p = prob_box.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = Inches(0.3)
    tf_p.margin_right = Inches(0.3)
    tf_p.margin_top = Inches(0.3)

    p = tf_p.paragraphs[0]
    p.text = "❌ The Problem"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)

    bullets_p = [
        "Traditional course creation takes days of manual content curation.",
        "Static online learning materials fail to adapt to individual student skill levels.",
        "Lack of structured practical labs, sample code, and hands-on capstone projects.",
        "Difficulty for learners to export customized curricula for offline study."
    ]
    for b in bullets_p:
        p_b = tf_p.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_before = Pt(12)

    sol_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(2.0), Inches(5.6), Inches(4.8))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = CARD_BG
    sol_box.line.color.rgb = ACCENT_CYAN
    sol_box.line.width = Pt(1.5)

    tf_s = sol_box.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = Inches(0.3)
    tf_s.margin_right = Inches(0.3)
    tf_s.margin_top = Inches(0.3)

    p = tf_s.paragraphs[0]
    p.text = "💡 CourseCraft AI Solution"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    bullets_s = [
        "Automated AI Course Synthesis Engine generates custom curricula in seconds.",
        "Supports dynamic filtering by Skill Level (Beginner to Advanced) & Duration.",
        "Generates multi-section theoretical lessons, industry context, pro-tips & code.",
        "Includes step-by-step terminal labs and evaluation-ready capstone projects.",
        "One-click export to Markdown (.md) and JSON for seamless learning."
    ]
    for b in bullets_s:
        p_b = tf_s.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_WHITE
        p_b.space_before = Pt(10)

    # SLIDE 3: Key Features & Core Modules
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3)
    add_header(slide3, "Key Features & Core Modules")

    features = [
        ("🔐 Auth System", "User registration, secure login with Werkzeug PBKDF2 hashing, session management, and admin role detection.", ACCENT_CYAN),
        ("🤖 AI Course Engine", "Generates full structured courses based on Topic, Category, Modules, Skill Level, and Custom Focus.", ACCENT_PURPLE),
        ("💻 Production Code", "Every module includes production-ready code with type annotations, defensive checks, and logging.", ACCENT_CYAN),
        ("🧪 Hands-On Labs", "Terminal step-by-step practical exercises & comprehensive final capstone project rubrics.", ACCENT_PURPLE),
        ("📄 Multi-Export", "Export generated courses instantly into formatted Markdown (.md) study guides or JSON data files.", ACCENT_CYAN),
        ("🎨 Glassmorphism UI", "Modern UI with dark mode, responsive cards, and an interactive 60fps HTML5 Canvas particle background.", ACCENT_PURPLE)
    ]

    for i, (title, desc, color) in enumerate(features):
        row = i // 3
        col = i % 3
        x = Inches(0.8 + col * 3.98)
        y = Inches(2.0 + row * 2.5)

        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.77), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.2)

        tf_f = box.text_frame
        tf_f.word_wrap = True
        tf_f.margin_left = Inches(0.2)
        tf_f.margin_right = Inches(0.2)
        tf_f.margin_top = Inches(0.2)

        p = tf_f.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        p_desc = tf_f.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(8)

    # SLIDE 4: Technology Stack
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4)
    add_header(slide4, "Technology Stack & Frameworks")

    tech_categories = [
        ("Backend & Logic", ["Python 3.14 (Core Runtime)", "Flask 3.0 (Web Framework & Routing)", "Werkzeug (Security & Password Hashing)", "Gunicorn (Production WSGI Server)"], ACCENT_CYAN),
        ("Database Architecture", ["Dual Database Layer", "SQLite3 (Local Development Storage)", "PostgreSQL (Cloud Production Database)", "psycopg2 (PostgreSQL Driver)"], ACCENT_PURPLE),
        ("Frontend & UI Design", ["HTML5 & Semantic Page Structure", "Vanilla CSS3 (Glassmorphism & Flexbox/Grid)", "JavaScript ES6+ (Interactive Logic & API)", "HTML5 Canvas (60fps Interactive Particles)"], ACCENT_CYAN),
        ("Deployment & Libraries", ["Render.com (Live Cloud Hosting)", "Marked.js (Markdown Rendering)", "FontAwesome 6 (Vector Icons)", "Render.yaml (Infrastructure as Code)"], ACCENT_PURPLE)
    ]

    for i, (title, items, color) in enumerate(tech_categories):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 5.96)
        y = Inches(2.0 + row * 2.5)

        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.76), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.2)

        tf_t = box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = Inches(0.25)
        tf_t.margin_right = Inches(0.25)
        tf_t.margin_top = Inches(0.2)

        p = tf_t.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = color

        for item in items:
            p_item = tf_t.add_paragraph()
            p_item.text = "• " + item
            p_item.font.size = Pt(12)
            p_item.font.color.rgb = TEXT_WHITE
            p_item.space_before = Pt(4)

    # SLIDE 5: System Architecture & Workflow
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5)
    add_header(slide5, "System Architecture & Processing Workflow")

    steps = [
        ("1. User Request", "User inputs Topic, Skill Level, Duration & Modules in Generator Wizard.", ACCENT_CYAN),
        ("2. Flask Routing", "App processes form request, verifies authentication session.", ACCENT_PURPLE),
        ("3. Synthesis Engine", "course_generator.py synthesizes lessons, labs, code & capstone project.", ACCENT_CYAN),
        ("4. Dual DB Save", "Database layer stores course JSON structured data linked to User ID.", ACCENT_PURPLE),
        ("5. Learning Studio", "Jinja2 renders interactive studio view with Markdown formatting.", ACCENT_CYAN)
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.8 + i * 2.4)
        y = Inches(2.2)

        box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.133), Inches(4.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        tf_st = box.text_frame
        tf_st.word_wrap = True
        tf_st.margin_left = Inches(0.15)
        tf_st.margin_right = Inches(0.15)
        tf_st.margin_top = Inches(0.25)

        p = tf_st.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color

        p_desc = tf_st.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(14)

    # SLIDE 6: Database Schema & Entity Design
    slide6 = prs.slides.add_slide(blank_layout)
    add_background(slide6)
    add_header(slide6, "Database Schema (SQLite3 / PostgreSQL Dual Model)")

    t1 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8))
    t1.fill.solid()
    t1.fill.fore_color.rgb = CARD_BG
    t1.line.color.rgb = ACCENT_CYAN
    t1.line.width = Pt(1.5)

    tf_t1 = t1.text_frame
    tf_t1.word_wrap = True
    tf_t1.margin_left = Inches(0.3)
    tf_t1.margin_right = Inches(0.3)
    tf_t1.margin_top = Inches(0.3)

    p = tf_t1.paragraphs[0]
    p.text = "👤 Users Table Schema"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    user_cols = [
        "id: INTEGER / SERIAL (Primary Key)",
        "username: VARCHAR(255) (Unique, Not Null)",
        "email: VARCHAR(255) (Unique, Not Null)",
        "password_hash: TEXT (PBKDF2 Hash)",
        "is_admin: INTEGER (0 = User, 1 = Admin)",
        "created_at: TIMESTAMP (Default Current)"
    ]
    for col in user_cols:
        p_c = tf_t1.add_paragraph()
        p_c.text = "🔹 " + col
        p_c.font.size = Pt(13)
        p_c.font.color.rgb = TEXT_WHITE
        p_c.space_before = Pt(8)

    t2 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(2.0), Inches(5.6), Inches(4.8))
    t2.fill.solid()
    t2.fill.fore_color.rgb = CARD_BG
    t2.line.color.rgb = ACCENT_PURPLE
    t2.line.width = Pt(1.5)

    tf_t2 = t2.text_frame
    tf_t2.word_wrap = True
    tf_t2.margin_left = Inches(0.3)
    tf_t2.margin_right = Inches(0.3)
    tf_t2.margin_top = Inches(0.3)

    p = tf_t2.paragraphs[0]
    p.text = "📚 Courses Table Schema"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE

    course_cols = [
        "id: INTEGER / SERIAL (Primary Key)",
        "user_id: INTEGER (Foreign Key -> users.id)",
        "title: VARCHAR(255) (Course Title)",
        "description: TEXT (Course Overview)",
        "category, level, duration: VARCHAR",
        "modules_data: TEXT (JSON Structured Data)",
        "created_at: TIMESTAMP (Creation Date)"
    ]
    for col in course_cols:
        p_c = tf_t2.add_paragraph()
        p_c.text = "🔹 " + col
        p_c.font.size = Pt(13)
        p_c.font.color.rgb = TEXT_WHITE
        p_c.space_before = Pt(8)

    # SLIDE 7: Key Project Files & Code Architecture
    slide7 = prs.slides.add_slide(blank_layout)
    add_background(slide7)
    add_header(slide7, "Core Project Modules & File Architecture")

    file_modules = [
        ("app.py", "Flask Application Controller", "Handles application setup, routing, authentication authorization check, user dashboard views, and export endpoints.", ACCENT_CYAN),
        ("course_generator.py", "AI Synthesis Engine", "Modular generation logic that dynamically synthesizes complete course JSON, lessons, code walkthroughs, labs, and capstones.", ACCENT_PURPLE),
        ("database.py", "Dual-Database Adapter", "Abstraction layer supporting both SQLite3 (local) and PostgreSQL (cloud). Handles password hashing & query adaptation.", ACCENT_CYAN),
        ("static/ & templates/", "Glassmorphism UI Engine", "Jinja2 templates (base, dashboard, studio) paired with custom CSS glassmorphism styles and 60fps Canvas particle network.", ACCENT_PURPLE)
    ]

    for i, (fname, frole, fdesc, color) in enumerate(file_modules):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 5.96)
        y = Inches(2.0 + row * 2.5)

        box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.76), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.2)

        tf_fm = box.text_frame
        tf_fm.word_wrap = True
        tf_fm.margin_left = Inches(0.25)
        tf_fm.margin_right = Inches(0.25)
        tf_fm.margin_top = Inches(0.2)

        p = tf_fm.paragraphs[0]
        p.text = f"📄 {fname}"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = color

        p_sub = tf_fm.add_paragraph()
        p_sub.text = frole
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.space_before = Pt(2)

        p_d = tf_fm.add_paragraph()
        p_d.text = fdesc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_WHITE
        p_d.space_before = Pt(6)

    # SLIDE 8: Deployment & Live Production
    slide8 = prs.slides.add_slide(blank_layout)
    add_background(slide8)
    add_header(slide8, "Production Cloud Deployment & Live App")

    dep_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8))
    dep_box.fill.solid()
    dep_box.fill.fore_color.rgb = CARD_BG
    dep_box.line.color.rgb = ACCENT_CYAN
    dep_box.line.width = Pt(1.5)

    tf_dp = dep_box.text_frame
    tf_dp.word_wrap = True
    tf_dp.margin_left = Inches(0.3)
    tf_dp.margin_right = Inches(0.3)
    tf_dp.margin_top = Inches(0.3)

    p = tf_dp.paragraphs[0]
    p.text = "☁️ Render Cloud Deployment"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    deploy_points = [
        "Infrastructure configured via render.yaml (Infrastructure as Code).",
        "WSGI Server: Gunicorn production process manager.",
        "Automatic database switching between local SQLite3 and production PostgreSQL.",
        "SSL/TLS Encryption (HTTPS) enabled out of the box."
    ]
    for pt in deploy_points:
        p_pt = tf_dp.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = TEXT_WHITE
        p_pt.space_before = Pt(12)

    demo_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(2.0), Inches(5.6), Inches(4.8))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = CARD_BG
    demo_box.line.color.rgb = ACCENT_PURPLE
    demo_box.line.width = Pt(1.5)

    tf_dm = demo_box.text_frame
    tf_dm.word_wrap = True
    tf_dm.margin_left = Inches(0.3)
    tf_dm.margin_right = Inches(0.3)
    tf_dm.margin_top = Inches(0.3)

    p = tf_dm.paragraphs[0]
    p.text = "🌐 Live Demo Information"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE

    p_url = tf_dm.add_paragraph()
    p_url.text = "Live Application URL:"
    p_url.font.size = Pt(13)
    p_url.font.color.rgb = TEXT_MUTED
    p_url.space_before = Pt(10)

    p_link = tf_dm.add_paragraph()
    p_link.text = "course-craft-ai-xlxe.onrender.com"
    p_link.font.size = Pt(14)
    p_link.font.bold = True
    p_link.font.color.rgb = ACCENT_CYAN
    p_link.space_before = Pt(4)

    p_gh = tf_dm.add_paragraph()
    p_gh.text = "\nGitHub Repository:"
    p_gh.font.size = Pt(13)
    p_gh.font.color.rgb = TEXT_MUTED

    p_ghlink = tf_dm.add_paragraph()
    p_ghlink.text = "github.com/Munim-Gupta/Course-Craft-Ai"
    p_ghlink.font.size = Pt(14)
    p_ghlink.font.bold = True
    p_ghlink.font.color.rgb = ACCENT_PURPLE
    p_ghlink.space_before = Pt(4)

    # SLIDE 9: Future Scope & Enhancements
    slide9 = prs.slides.add_slide(blank_layout)
    add_background(slide9)
    add_header(slide9, "Future Enhancements & Roadmap")

    futures = [
        ("🤖 LLM API Integration", "Connect with OpenAI GPT-4 / Google Gemini API for real-time generative AI course content expansion.", ACCENT_CYAN),
        ("📊 Student Analytics", "Track course completion progress, interactive quiz scores, and student learning metrics.", ACCENT_PURPLE),
        ("📑 PDF & EPUB Export", "Extend export engine to support stylized PDF e-books and EPUB digital readers.", ACCENT_CYAN),
        ("🏆 Quiz & Certification", "Auto-generate module quizzes and automated certificates of achievement upon completion.", ACCENT_PURPLE)
    ]

    for i, (title, desc, color) in enumerate(futures):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 5.96)
        y = Inches(2.0 + row * 2.5)

        box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.76), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        box.line.width = Pt(1.2)

        tf_fu = box.text_frame
        tf_fu.word_wrap = True
        tf_fu.margin_left = Inches(0.25)
        tf_fu.margin_right = Inches(0.25)
        tf_fu.margin_top = Inches(0.2)

        p = tf_fu.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = color

        p_d = tf_fu.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_WHITE
        p_d.space_before = Pt(8)

    # SLIDE 10: Conclusion & Q&A
    slide10 = prs.slides.add_slide(blank_layout)
    add_background(slide10)

    card10 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card10.fill.solid()
    card10.fill.fore_color.rgb = CARD_BG
    card10.line.color.rgb = ACCENT_CYAN
    card10.line.width = Pt(1.5)

    tb10 = slide10.shapes.add_textbox(Inches(2.0), Inches(1.8), Inches(9.333), Inches(3.8))
    tf10 = tb10.text_frame
    tf10.word_wrap = True

    p = tf10.paragraphs[0]
    p.text = "Thank You! 🙏"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.CENTER

    p2 = tf10.add_paragraph()
    p2.text = "Questions & Answers Session"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf10.add_paragraph()
    p3.text = "\nCourseCraft AI — SYIT Project Presentation"
    p3.font.size = Pt(16)
    p3.font.color.rgb = ACCENT_PURPLE
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf10.add_paragraph()
    p4.text = "\nStudent: Munim Gupta  |  Live App: https://course-craft-ai-xlxe.onrender.com"
    p4.font.size = Pt(14)
    p4.font.color.rgb = TEXT_MUTED
    p4.alignment = PP_ALIGN.CENTER

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CourseCraft_AI_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
